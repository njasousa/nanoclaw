from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
SAP_BLUE = RGBColor(0x00, 0x6F, 0xCA)       # SAP Blue
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)      # Dark Blue
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)    # Light Blue background
GOLD = RGBColor(0xF0, 0xA5, 0x00)           # Gold/Amber accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x86, 0x48)
ORANGE = RGBColor(0xE8, 0x6A, 0x10)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=12, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_slide(slide, title_text, bullets, left_col=0.4, top_start=1.6, font_size=13):
    """Add a title + bulleted list to a slide."""
    # Title bar already added by caller; just add bullets
    y = top_start
    for bullet in bullets:
        indent = bullet.get("indent", 0)
        text = bullet["text"]
        size = bullet.get("size", font_size)
        bold = bullet.get("bold", False)
        color = bullet.get("color", DARK_GRAY)
        prefix = "  " * indent + ("• " if indent == 0 else "  – ")
        add_textbox(slide, prefix + text, left_col + indent * 0.2, y, 12.3 - indent * 0.2, 0.45,
                    font_size=size, bold=bold, color=color)
        y += bullet.get("spacing", 0.38)
    return y


def slide_header(slide, title, subtitle=None, bar_color=DARK_BLUE):
    """Add top header bar + title."""
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=bar_color)
    add_textbox(slide, title, 0.3, 0.12, 12.5, 0.7,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.75, 12.5, 0.4,
                    font_size=13, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
    # Bottom accent line
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=SAP_BLUE)

def slide_footer(slide, text="CX Solutions | SAP Activate Framework"):
    add_textbox(slide, text, 0.3, 7.22, 12.5, 0.25,
                font_size=8, color=WHITE, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, fill_color=GOLD)      # top stripe
add_rect(slide, 0, 7.42, 13.33, 0.08, fill_color=GOLD)   # bottom stripe
add_rect(slide, 0, 2.4, 13.33, 2.7, fill_color=SAP_BLUE) # center band

add_textbox(slide, "Customer Experience (CX)", 0.5, 2.55, 12.3, 0.8,
            font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Da Necessidade ao ROI: Uma Abordagem Estruturada com SAP Activate", 0.5, 3.35, 12.3, 0.6,
            font_size=18, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

add_textbox(slide, "Transformação Digital • E-Commerce • Resultados Mensuráveis", 0.5, 6.5, 12.3, 0.4,
            font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "Agenda", bar_color=DARK_BLUE)
slide_footer(slide)

items = [
    ("01", "A Necessidade de Soluções CX nas Empresas"),
    ("02", "SAP Activate Framework — Fases, Inputs & Outputs"),
    ("03", "Milestones Críticos do Projeto"),
    ("04", "KPIs que os CxOs Monitorizam"),
    ("05", "Impacto nas Métricas após Go-Live"),
    ("06", "ROI, Payback Period e Benchmarks"),
]
for i, (num, label) in enumerate(items):
    y = 1.45 + i * 0.87
    add_rect(slide, 0.4, y, 1.0, 0.65, fill_color=SAP_BLUE)
    add_textbox(slide, num, 0.4, y + 0.08, 1.0, 0.5, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, y, 11.0, 0.65, fill_color=WHITE)
    add_textbox(slide, label, 1.65, y + 0.1, 10.7, 0.5, font_size=15, bold=False, color=DARK_BLUE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — NECESSIDADE CX
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "A Necessidade de Soluções CX", "Por que razão as empresas investem em plataformas de Comércio Eletrónico e CX?")
slide_footer(slide)

# Left column
add_rect(slide, 0.3, 1.35, 6.0, 5.7, fill_color=WHITE)
add_textbox(slide, "🌐  Contexto de Mercado", 0.5, 1.45, 5.6, 0.4, font_size=13, bold=True, color=DARK_BLUE)
bullets_left = [
    {"text": "89% dos consumidores mudaram para um concorrente após uma má experiência (Gartner)", "indent": 0},
    {"text": "O mercado global de e-commerce ultrapassa os $6,3 triliões em 2024 (Statista)", "indent": 0},
    {"text": "As empresas com CX de excelência crescem 4–8× mais do que a média da indústria (Bain)", "indent": 0},
    {"text": "70% das decisões de compra B2B são influenciadas pela experiência digital (Forrester)", "indent": 0},
    {"text": "A omnicanalidade deixou de ser diferenciador — é requisito mínimo", "indent": 0},
]
y = 1.9
for b in bullets_left:
    add_textbox(slide, "• " + b["text"], 0.5, y, 5.6, 0.65, font_size=11.5, color=DARK_GRAY)
    y += 0.68

# Right column
add_rect(slide, 6.7, 1.35, 6.3, 5.7, fill_color=WHITE)
add_textbox(slide, "⚡  Desafios que Motivam o Investimento", 6.9, 1.45, 5.9, 0.4, font_size=13, bold=True, color=DARK_BLUE)
challenges = [
    ("Silos de dados", "Experiência fragmentada entre canais online e offline"),
    ("Processos manuais", "Catálogos, preços e stock geridos manualmente → erros e lentidão"),
    ("Falta de personalização", "Clientes esperam recomendações e jornadas personalizadas"),
    ("Time-to-market lento", "Promoções e novos produtos demoram semanas a lançar"),
    ("Visibilidade limitada", "Sem analytics em tempo real sobre comportamento e conversão"),
]
y = 1.9
for title, desc in challenges:
    add_rect(slide, 6.8, y, 0.08, 0.5, fill_color=SAP_BLUE)
    add_textbox(slide, title, 7.0, y, 5.8, 0.22, font_size=11.5, bold=True, color=DARK_BLUE)
    add_textbox(slide, desc, 7.0, y + 0.22, 5.8, 0.28, font_size=10.5, color=MID_GRAY)
    y += 0.72

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — SAP ACTIVATE OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "SAP Activate Framework", "Metodologia ágil e estruturada para implementações SAP — adaptada a projetos CX")
slide_footer(slide)

phases = [
    ("Discover", SAP_BLUE),
    ("Prepare", RGBColor(0x00, 0x8A, 0xC9)),
    ("Explore", RGBColor(0x00, 0xA3, 0xA3)),
    ("Realize", RGBColor(0x2E, 0x86, 0x48)),
    ("Deploy", GOLD),
    ("Run", ORANGE),
]
w = 13.33 / 6
for i, (name, color) in enumerate(phases):
    x = i * w
    add_rect(slide, x + 0.05, 1.4, w - 0.1, 0.7, fill_color=color)
    add_textbox(slide, name, x + 0.05, 1.4, w - 0.1, 0.7,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow connector (except last)
    if i < 5:
        add_rect(slide, x + w - 0.1, 1.62, 0.2, 0.26, fill_color=color)

add_textbox(slide, "Fundamentos do SAP Activate:", 0.4, 2.3, 12.5, 0.35, font_size=13, bold=True, color=DARK_BLUE)

cols = [
    ("Abordagem Fit-to-Standard", "Adaptar processos de negócio à solução, minimizando customizações desnecessárias."),
    ("Entregas Iterativas (Sprints)", "Ciclos curtos com demos regulares ao cliente para validação contínua."),
    ("Acelerators & Best Practices", "Conteúdos pré-configurados (iBP, roadmaps) que reduzem tempo de implementação."),
    ("Qualidade em cada Gate", "Quality Gates formais garantem que cada fase está completa antes de avançar."),
]
for i, (title, desc) in enumerate(cols):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 2.75, 3.0, 4.3, fill_color=WHITE)
    add_rect(slide, x, 2.75, 3.0, 0.08, fill_color=SAP_BLUE)
    add_textbox(slide, title, x + 0.1, 2.85, 2.8, 0.5, font_size=11.5, bold=True, color=DARK_BLUE)
    add_textbox(slide, desc, x + 0.1, 3.35, 2.8, 1.5, font_size=11, color=MID_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — PHASE: DISCOVER
# ─────────────────────────────────────────────────────────────────────────────
def phase_slide(prs, phase_name, phase_color, phase_num, duration, objective,
                inputs, outputs, activities, milestone=None):
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)

    # Header with phase color
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=phase_color)
    add_textbox(slide, f"Fase {phase_num}: {phase_name}", 0.3, 0.08, 10.0, 0.6,
                font_size=26, bold=True, color=WHITE)
    add_textbox(slide, f"⏱ Duração típica: {duration}", 0.3, 0.72, 6.0, 0.38,
                font_size=12, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=SAP_BLUE)
    slide_footer(slide)

    # Objective banner
    add_rect(slide, 0.3, 1.3, 12.73, 0.5, fill_color=RGBColor(0xE0, 0xF0, 0xFF))
    add_textbox(slide, "🎯  " + objective, 0.45, 1.32, 12.5, 0.44, font_size=12, bold=False, color=DARK_BLUE)

    # Three columns: Activities | Inputs | Outputs
    # Activities
    add_rect(slide, 0.3, 1.9, 4.1, 5.1, fill_color=WHITE)
    add_rect(slide, 0.3, 1.9, 4.1, 0.42, fill_color=phase_color)
    add_textbox(slide, "Atividades Principais", 0.4, 1.93, 3.9, 0.36, font_size=12, bold=True, color=WHITE)
    y = 2.4
    for act in activities:
        add_textbox(slide, "▸ " + act, 0.4, y, 3.9, 0.55, font_size=10.5, color=DARK_GRAY)
        y += 0.52

    # Inputs
    add_rect(slide, 4.6, 1.9, 4.0, 5.1, fill_color=WHITE)
    add_rect(slide, 4.6, 1.9, 4.0, 0.42, fill_color=DARK_BLUE)
    add_textbox(slide, "Inputs (Entradas)", 4.7, 1.93, 3.8, 0.36, font_size=12, bold=True, color=WHITE)
    y = 2.4
    for inp in inputs:
        add_textbox(slide, "📥 " + inp, 4.7, y, 3.8, 0.55, font_size=10.5, color=DARK_GRAY)
        y += 0.52

    # Outputs
    add_rect(slide, 8.8, 1.9, 4.23, 5.1, fill_color=WHITE)
    add_rect(slide, 8.8, 1.9, 4.23, 0.42, fill_color=GREEN)
    add_textbox(slide, "Outputs (Saídas)", 8.9, 1.93, 4.03, 0.36, font_size=12, bold=True, color=WHITE)
    y = 2.4
    for out in outputs:
        add_textbox(slide, "📤 " + out, 8.9, y, 4.03, 0.55, font_size=10.5, color=DARK_GRAY)
        y += 0.52

    if milestone:
        add_rect(slide, 0.3, 6.85, 12.73, 0.32, fill_color=GOLD)
        add_textbox(slide, "🏁 MILESTONE: " + milestone, 0.45, 6.87, 12.5, 0.28,
                    font_size=11, bold=True, color=DARK_BLUE)

phase_slide(
    prs, "Discover", SAP_BLUE, "1", "2–4 semanas",
    "Validar a visão de negócio, identificar o âmbito da solução CX e construir o business case.",
    inputs=["Estratégia de negócio e objetivos corporativos", "Análise de mercado e benchmarks",
            "Identificação de pain points atuais", "Orçamento preliminar e restrições",
            "Stakeholders e sponsors identificados"],
    outputs=["Business Case validado", "Visão de alto nível da solução CX",
             "Âmbito preliminar (scope statement)", "Análise de fit SAP CX (Commerce, Sales, Service)",
             "Recomendação de roadmap de implementação"],
    activities=["Workshops de descoberta com C-Suite", "Demo da solução SAP CX",
                "Análise de processos AS-IS", "Avaliação de maturidade digital",
                "Definição de critérios de sucesso"],
    milestone="Business Case Aprovado & Decisão de Avançar (Go/No-Go)"
)

phase_slide(
    prs, "Prepare", RGBColor(0x00, 0x8A, 0xC9), "2", "4–6 semanas",
    "Constituir a equipa de projeto, estabelecer o ambiente técnico e definir o plano detalhado.",
    inputs=["Business Case aprovado", "Contrato e SoW assinados",
            "Recursos de negócio e IT identificados", "Landscape tecnológico atual (ERP, CRM, PIM)",
            "Requisitos de integração de alto nível"],
    outputs=["Project Charter e plano de projeto detalhado", "Equipa de projeto constituída (roles & responsibilities)",
             "Ambientes técnicos provisionados (Dev/QA/Prod)", "Backlog inicial de épicos e user stories",
             "Plano de gestão de mudança (Change Management Plan)", "Definition of Done & quality standards"],
    activities=["Kick-off do projeto com todos os stakeholders", "Setup dos ambientes SAP CX",
                "Formação inicial da equipa em SAP Activate", "Elaboração do project plan (MS Project / Jira)",
                "Definição do modelo de governança"],
    milestone="Project Kick-off Realizado & Ambientes Técnicos Prontos"
)

phase_slide(
    prs, "Explore", RGBColor(0x00, 0xA3, 0xA3), "3", "6–10 semanas",
    "Validar os processos de negócio contra a solução padrão SAP CX (Fit-to-Standard Workshops).",
    inputs=["Backlog de user stories priorizado", "Catálogo de processos AS-IS",
            "Dados de referência (produtos, clientes, preços)", "Requisitos de integração detalhados",
            "Restrições técnicas e de segurança"],
    outputs=["Documento de Design de Solução (SDD / BPD)", "Gap List com análise Fit/Gap/Enhancement",
             "Protótipo / PoC de cenários críticos", "Especificações de integração (API design)",
             "Backlog refinado com estimativas", "Acceptance Criteria validados pelo negócio"],
    activities=["Fit-to-Standard Workshops por módulo (Commerce, OMS, Search)", "Demonstrações do sistema configurado",
                "Registo de gaps e decisões de design", "Prototipagem de jornadas do cliente",
                "Revisão de integrações com ERP/CRM/PIM"],
    milestone="Design da Solução Aprovado — Quality Gate Fase Explore"
)

phase_slide(
    prs, "Realize", GREEN, "4", "12–20 semanas",
    "Construir, configurar, integrar e testar a solução CX de acordo com o design aprovado.",
    inputs=["Solution Design Document aprovado", "Backlog refinado e estimado",
            "Dados mestres preparados (catálogo, clientes)", "Especificações de integração finalizadas",
            "Plano de testes (Test Plan)"],
    outputs=["Sistema configurado e customizações desenvolvidas", "Integrações implementadas e testadas (Unit/Integration Tests)",
             "UAT (User Acceptance Testing) executado e aprovado", "Dados migrados e validados",
             "Documentação técnica e funcional", "Plano de formação de utilizadores finais"],
    activities=["Sprints de configuração e desenvolvimento (2 semanas cada)", "Demos de sprint ao Product Owner",
                "Testes de integração com ERP, CRM, PIM, pagamentos", "Testes de performance e carga",
                "Formação de key users e super users", "Preparação da estratégia de cutover"],
    milestone="UAT Assinado & System Integration Tests Aprovados"
)

phase_slide(
    prs, "Deploy", GOLD, "5", "3–5 semanas",
    "Preparar e executar o go-live da plataforma CX em produção com mínimo impacto operacional.",
    inputs=["UAT aprovado (sign-off)", "Plano de cutover detalhado",
            "Dados migrados e validados em produção", "Utilizadores formados e certificados",
            "Runbook de go-live e plano de rollback"],
    outputs=["Sistema em produção (Go-Live)", "Hipercare plan ativado",
             "Suporte de 1ª linha operacional", "Relatório de go-live",
             "Documentação de operações entregue", "Handover para equipa de Run"],
    activities=["Ensaio de cutover (Cutover Rehearsal)", "Freeze de alterações ao sistema legacy",
                "Migração final de dados para produção", "Smoke tests e validação go-live",
                "Comunicação a clientes e parceiros", "Ativação faseada (soft launch → full launch)"],
    milestone="GO-LIVE — Plataforma CX em Produção"
)

phase_slide(
    prs, "Run", ORANGE, "6", "Contínuo (3–6 meses de hipercare + ongoing)",
    "Estabilizar a solução em produção, otimizar com base em dados reais e evoluir a plataforma.",
    inputs=["Sistema em produção", "Feedback de utilizadores e clientes",
            "Tickets de suporte e incidentes", "KPIs e dashboards de performance",
            "Backlog de melhorias (Enhancement Backlog)"],
    outputs=["SLA de suporte cumprido", "Relatórios de performance (disponibilidade, conversão, NPS)",
             "Melhorias iterativas implementadas", "Plano de evolução da plataforma (Roadmap v2)",
             "Lições aprendidas documentadas", "Business Value Realization Report"],
    activities=["Monitorização contínua (uptime, performance, segurança)", "Resolução de incidentes e bugs",
                "Otimização SEO, UX e conversão com base em analytics", "Formação de novos utilizadores",
                "Revisões periódicas de KPIs com o negócio", "Planeamento de releases e evoluções"],
    milestone="Fim do Período de Hipercare & Transição para BAU (Business As Usual)"
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — MILESTONES TIMELINE
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "Milestones Críticos do Projeto CX", "Do Business Case ao Business As Usual")
slide_footer(slide)

milestones = [
    ("M1", "Business Case\nAprovado", SAP_BLUE, "Discover", "Sem. 4"),
    ("M2", "Kick-off &\nAmbientes Prontos", RGBColor(0x00, 0x8A, 0xC9), "Prepare", "Sem. 8"),
    ("M3", "Design da\nSolução Aprovado", RGBColor(0x00, 0xA3, 0xA3), "Explore", "Sem. 18"),
    ("M4", "UAT Sign-off &\nSIT Aprovado", GREEN, "Realize", "Sem. 38"),
    ("M5", "GO-LIVE", GOLD, "Deploy", "Sem. 43"),
    ("M6", "Fim Hipercare\n& BAU", ORANGE, "Run", "Sem. 55"),
]

# Timeline bar
add_rect(slide, 0.5, 2.8, 12.33, 0.12, fill_color=DARK_BLUE)

for i, (code, label, color, phase, week) in enumerate(milestones):
    x = 0.5 + i * 2.05
    # Diamond marker
    add_rect(slide, x + 0.65, 2.55, 0.35, 0.35, fill_color=color)
    # Code box
    add_rect(slide, x + 0.55, 1.5, 0.55, 0.55, fill_color=color)
    add_textbox(slide, code, x + 0.55, 1.52, 0.55, 0.5, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Label below line
    add_rect(slide, x + 0.2, 3.1, 1.6, 1.1, fill_color=WHITE)
    add_rect(slide, x + 0.2, 3.1, 1.6, 0.05, fill_color=color)
    add_textbox(slide, label, x + 0.22, 3.18, 1.56, 0.65, font_size=9.5, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, week, x + 0.22, 3.85, 1.56, 0.3, font_size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Phase label above
    add_textbox(slide, phase, x + 0.2, 2.12, 1.6, 0.3, font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Key notes
add_rect(slide, 0.3, 4.35, 12.73, 2.7, fill_color=WHITE)
add_textbox(slide, "Pontos-Chave sobre Milestones", 0.5, 4.42, 12.3, 0.35, font_size=13, bold=True, color=DARK_BLUE)
notes = [
    "Quality Gates formais: cada milestone requer aprovação formal dos stakeholders antes de avançar para a fase seguinte.",
    "Tolerância a desvios: atrasos no M3 (Design) propagam-se a todas as fases — é o milestone mais crítico para o prazo global.",
    "Go-Live faseado: recomendado lançar por mercado, canal ou linha de produto para mitigar risco operacional.",
    "Hipercare (M5→M6): período intensivo de suporte pós go-live, tipicamente 6–12 semanas com equipa de projeto disponível 24/7.",
    "Duração total típica: projetos CX de média dimensão levam 10–14 meses do kick-off ao fim de hipercare.",
]
y = 4.85
for note in notes:
    add_textbox(slide, "▸ " + note, 0.5, y, 12.4, 0.38, font_size=11, color=DARK_GRAY)
    y += 0.4

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — KPIs CxOs
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "KPIs Monitorizados pelos CxOs", "Métricas que provam o valor de uma plataforma CX")
slide_footer(slide)

cxos = [
    ("CEO", SAP_BLUE, [
        "Net Promoter Score (NPS)",
        "Crescimento de Receita Online",
        "Market Share Digital",
        "Customer Lifetime Value (CLV)",
        "Satisfação Global do Cliente (CSAT)",
    ]),
    ("CFO", DARK_BLUE, [
        "ROI da plataforma CX",
        "Custo de Aquisição de Cliente (CAC)",
        "Revenue per Visitor (RPV)",
        "Margem por canal digital",
        "Payback Period do investimento",
    ]),
    ("CIO / CTO", RGBColor(0x00, 0xA3, 0xA3), [
        "Disponibilidade da plataforma (Uptime)",
        "Tempo de resposta / Page Load Time",
        "Nº de integrações activas e estáveis",
        "Incidentes de segurança / DORA Metrics",
        "Time-to-Deploy de novas funcionalidades",
    ]),
    ("CMO", GREEN, [
        "Conversion Rate (CVR)",
        "Taxa de abandono do carrinho",
        "Custo por Clique e ROAS",
        "Tráfego orgânico (SEO) e engagement",
        "Taxa de retenção e compras repetidas",
    ]),
]
for i, (role, color, kpis) in enumerate(cxos):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 1.4, 3.05, 5.65, fill_color=WHITE)
    add_rect(slide, x, 1.4, 3.05, 0.5, fill_color=color)
    add_textbox(slide, role, x, 1.42, 3.05, 0.46, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y = 2.05
    for kpi in kpis:
        add_rect(slide, x + 0.12, y, 0.06, 0.24, fill_color=color)
        add_textbox(slide, kpi, x + 0.25, y, 2.7, 0.42, font_size=10.5, color=DARK_GRAY)
        y += 0.88

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — IMPACTO NAS MÉTRICAS
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "Impacto nas Métricas após Implementação CX",
             "Benchmarks baseados em estudos da Forrester, McKinsey, SAP e Gartner")
slide_footer(slide)

metrics = [
    ("📈 Receita Online", "+25–40%", "Aumento médio no primeiro ano após go-live", GREEN),
    ("🛒 Conversion Rate", "+15–30%", "Melhoria com UX optimizado e personalização", GREEN),
    ("🔁 Taxa de Retenção", "+10–20%", "Clientes que voltam a comprar (repeat purchase)", GREEN),
    ("🎯 NPS", "+15–25 pontos", "Incremento no Net Promoter Score", GREEN),
    ("⚡ Time-to-Market", "-40–60%", "Redução no tempo de lançamento de produtos/campanhas", SAP_BLUE),
    ("💰 Custo Operacional", "-20–35%", "Automatização de processos manuais (orders, fulfillment)", SAP_BLUE),
    ("📊 CAC (Custo Aquisição)", "-15–25%", "Melhor targeting e personalização reduzem custo", SAP_BLUE),
    ("🏎 Page Load Time", "-50–70%", "Plataformas modernas com CDN e optimização", SAP_BLUE),
]

for i, (metric, value, desc, color) in enumerate(metrics):
    row = i // 4
    col = i % 4
    x = 0.3 + col * 3.2
    y = 1.4 + row * 2.85
    add_rect(slide, x, y, 3.05, 2.6, fill_color=WHITE)
    add_rect(slide, x, y, 3.05, 0.06, fill_color=color)
    add_textbox(slide, metric, x + 0.1, y + 0.12, 2.85, 0.4, font_size=11, bold=True, color=DARK_BLUE)
    add_textbox(slide, value, x + 0.1, y + 0.58, 2.85, 0.65, font_size=22, bold=True, color=color)
    add_textbox(slide, desc, x + 0.1, y + 1.3, 2.85, 0.8, font_size=9.5, color=MID_GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — ROI & PAYBACK
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
slide_header(slide, "ROI e Payback Period", "Quanto tempo demora o retorno e qual o ROI médio de uma plataforma CX?")
slide_footer(slide)

# Big number boxes
stats = [
    ("14–24 meses", "Payback Period Médio", "Tempo até recuperar o investimento inicial", SAP_BLUE),
    ("200–400%", "ROI a 3 Anos", "Retorno sobre o investimento num horizonte de 3 anos", GREEN),
    ("$3.5M", "Benefício Médio Anual", "Para empresas mid-market (Forrester TEI Studies)", GOLD),
    ("18 meses", "Time-to-Value", "Tempo médio para impacto significativo nas métricas", ORANGE),
]
for i, (value, label, desc, color) in enumerate(stats):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 1.4, 3.05, 2.5, fill_color=color)
    add_textbox(slide, value, x + 0.1, 1.5, 2.85, 0.8, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x + 0.1, 2.3, 2.85, 0.4, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x + 0.1, 2.78, 2.85, 0.55, font_size=9.5, color=RGBColor(0xEE, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

# Fatores que influenciam o ROI
add_rect(slide, 0.3, 4.05, 6.0, 3.1, fill_color=WHITE)
add_rect(slide, 0.3, 4.05, 6.0, 0.42, fill_color=DARK_BLUE)
add_textbox(slide, "Fatores que Aceleram o ROI", 0.4, 4.07, 5.8, 0.38, font_size=12, bold=True, color=WHITE)
positive = [
    "Adoção rápida pelos utilizadores (Change Management eficaz)",
    "Integração nativa com ERP/CRM (menos customizações)",
    "Fit-to-Standard elevado (menos desenvolvimento à medida)",
    "Maturidade dos dados mestres (catálogo, preços, stock)",
    "Estratégia omnicanal clara desde o início",
]
y = 4.58
for p in positive:
    add_textbox(slide, "✓  " + p, 0.45, y, 5.7, 0.38, font_size=10.5, color=DARK_GRAY)
    y += 0.42

add_rect(slide, 6.7, 4.05, 6.33, 3.1, fill_color=WHITE)
add_rect(slide, 6.7, 4.05, 6.33, 0.42, fill_color=ORANGE)
add_textbox(slide, "Fatores que Diminuem / Atrasam o ROI", 6.8, 4.07, 6.13, 0.38, font_size=12, bold=True, color=WHITE)
negative = [
    "Scope creep e requisitos mal definidos",
    "Dados mestres de fraca qualidade (catálogo incompleto)",
    "Resistência à mudança interna (adoção baixa)",
    "Integrações complexas com sistemas legacy",
    "Falta de recursos de negócio dedicados ao projeto",
]
y = 4.58
for n in negative:
    add_textbox(slide, "✗  " + n, 6.8, y, 6.1, 0.38, font_size=10.5, color=DARK_GRAY)
    y += 0.42

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — CONCLUSION
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, fill_color=GOLD)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill_color=GOLD)

add_textbox(slide, "Conclusões Chave", 0.5, 0.3, 12.3, 0.7,
            font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

conclusions = [
    ("🎯 Necessidade Estratégica",
     "Investir em CX não é opcional — é fator de sobrevivência competitiva. Empresas sem plataforma digital estruturada perdem quota de mercado para concorrentes nativamente digitais."),
    ("🏗 Metodologia SAP Activate",
     "As 6 fases (Discover → Run) garantem estrutura, qualidade e previsibilidade. Inputs e outputs bem definidos em cada fase reduzem risco e aumentam a taxa de sucesso do projeto."),
    ("📊 ROI Comprovado",
     "Com ROI médio de 200–400% a 3 anos e payback entre 14–24 meses, o investimento em CX tem retorno claro e mensurável — desde que executado com rigor e adoção adequada."),
    ("🔑 Fatores Críticos de Sucesso",
     "Patrocínio executivo ativo, qualidade de dados, change management e fit-to-standard são os 4 pilares que determinam se um projeto CX atinge o seu potencial de valor."),
]
y = 1.2
for title, desc in conclusions:
    add_rect(slide, 0.5, y, 12.33, 1.25, fill_color=SAP_BLUE)
    add_rect(slide, 0.5, y, 0.08, 1.25, fill_color=GOLD)
    add_textbox(slide, title, 0.7, y + 0.08, 11.9, 0.35, font_size=13, bold=True, color=GOLD)
    add_textbox(slide, desc, 0.7, y + 0.45, 11.9, 0.65, font_size=11.5, color=RGBColor(0xDD, 0xEE, 0xFF))
    y += 1.4

add_textbox(slide, "Obrigado  •  Q&A", 0.5, 6.85, 12.3, 0.4,
            font_size=14, bold=True, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
output_path = "/Users/cane/Desktop/CX_SAP_Activate_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
